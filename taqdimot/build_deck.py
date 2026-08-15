"""UniAgent — NEXUS30 pitch deck generator.

Builds `UniAgent_pitch.pptx` (16:9, 13 slides) from the screenshots in
`screenshots/`. Requires `python-pptx` (not a project dependency — install it
into any throwaway environment: `pip install python-pptx`).

    python build_deck.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
SHOTS = HERE / "screenshots"
OUT = HERE / "UniAgent_pitch.pptx"

# --------------------------------------------------------------- palette
BG = RGBColor(0x0B, 0x10, 0x20)
CARD = RGBColor(0x14, 0x1D, 0x38)
CARD2 = RGBColor(0x1B, 0x26, 0x47)
LINE = RGBColor(0x2A, 0x37, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9D, 0xB0, 0xDC)
DIM = RGBColor(0x6D, 0x7E, 0xAC)
BLUE = RGBColor(0x5B, 0x8C, 0xFF)
MINT = RGBColor(0x2F, 0xD4, 0xA8)
AMBER = RGBColor(0xFF, 0xC2, 0x4B)
RED = RGBColor(0xFF, 0x6B, 0x6B)

FONT = "Calibri"
TOTAL = 13
ML = 0.70          # left margin
CW = 11.93         # content width

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# --------------------------------------------------------------- helpers
def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _paint(bg, BG)
    return s


def _paint(shape, fill, line=None, width=1.0):
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(width)


def card(slide, l, t, w, h, fill=CARD, line=None, radius=0.045):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    try:
        sh.adjustments[0] = radius
    except (IndexError, KeyError):
        pass
    _paint(sh, fill, line)
    return sh


def dot(slide, l, t, size, color, shape=MSO_SHAPE.OVAL):
    sh = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(size), Inches(size))
    _paint(sh, color)
    return sh


def text(slide, l, t, w, h, runs, size=12, color=MUTED, bold=False, align=PP_ALIGN.LEFT,
         spacing=1.0, anchor=MSO_ANCHOR.TOP, space_after=0, tracking=None):
    """`runs` is a str, a list of str (one paragraph each), or a list of
    (str, {overrides}) tuples for per-run styling inside one paragraph."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    paras = runs if isinstance(runs, list) else [runs]
    for i, item in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        pieces = item if isinstance(item, list) else [item]
        for piece in pieces:
            txt, over = piece if isinstance(piece, tuple) else (piece, {})
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = FONT
            f.size = Pt(over.get("size", size))
            f.bold = over.get("bold", bold)
            f.color.rgb = over.get("color", color)
            tr = over.get("tracking", tracking)
            if tr:
                f._rPr.set("spc", str(int(tr * 100)))
    return box


def eyebrow(slide, num, label, color=BLUE):
    text(slide, ML, 0.42, 8.0, 0.30, [[(f"{num:02d} — ", {"color": color}), (label.upper(), {"color": MUTED})]],
         size=12, bold=True, tracking=1.6)


def headline(slide, txt, size=32, color=WHITE, top=0.78, h=0.95, w=CW):
    text(slide, ML, top, w, h, txt, size=size, bold=True, color=color, spacing=1.05)


def footer(slide, index):
    text(slide, ML, 6.95, 4.0, 0.30, "UNIAGENT · NEXUS30", size=9, color=DIM, bold=True, tracking=1.2)
    text(slide, 8.63, 6.95, 4.0, 0.30, f"{index:02d} / {TOTAL}", size=9, color=DIM, bold=True,
         align=PP_ALIGN.RIGHT, tracking=1.2)


def source(slide, txt, top=6.55, w=11.93):
    text(slide, ML, top, w, 0.30, txt, size=9, color=DIM)


def shot(slide, name, l, t, w, crop=(0.0, 0.0, 0.0, 0.0), frame=True):
    """crop = (left, top, right, bottom) fractions. Height is derived so the
    visible region keeps the source pixel aspect (2880x1800)."""
    cl, ct, cr, cb = crop
    h = w * (1800 * (1 - ct - cb)) / (2880 * (1 - cl - cr))
    pic = slide.shapes.add_picture(str(SHOTS / name), Inches(l), Inches(t), Inches(w), Inches(h))
    pic.crop_left, pic.crop_top, pic.crop_right, pic.crop_bottom = cl, ct, cr, cb
    pic.shadow.inherit = False
    if frame:
        fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        _paint(fr, None, LINE, 1.0)
    return h


# =====================================================================
# 01 — title
# =====================================================================
s = new_slide()
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(2.10), Inches(1.05), Inches(0.075))
_paint(accent, BLUE)

text(s, ML, 2.45, 9.0, 1.30, "UniAgent", size=80, bold=True, color=WHITE)
text(s, ML, 3.78, 9.6, 0.55, "Universitetning kundalik ishi — bitta AI ish maydonida.",
     size=24, color=MUTED)
text(s, ML, 4.32, 9.6, 0.55, "Har rol o‘z doirasida. Har javob manba bilan.", size=24, bold=True, color=BLUE)

card(s, ML, 5.55, 5.30, 0.62, CARD, LINE)
text(s, 1.00, 5.72, 4.8, 0.30, "NEXUS30 · EdTech treki · hamkor: Yandex", size=12, bold=True, color=MUTED)
text(s, 8.63, 5.72, 4.0, 0.30, "Mannonov Sarabek · Avgust 2026", size=12, color=DIM, align=PP_ALIGN.RIGHT)

# =====================================================================
# 02 — background
# =====================================================================
s = new_slide()
eyebrow(s, 1, "Background")
headline(s, "O‘zbekiston oliy ta’limi 10 yilda olti barobar kengaydi")

card(s, ML, 1.90, CW, 1.62, CARD2)
text(s, 1.15, 2.10, 3.2, 1.10, "1.5 mln", size=64, bold=True, color=WHITE)
text(s, 4.55, 2.22, 7.6, 0.45, "talaba — 2025/26 o‘quv yilida", size=20, bold=True, color=WHITE)
text(s, 4.55, 2.72, 7.6, 0.65,
     "2015/16 da 264 ming edi. Har biri kundalik savoli, hujjati va arizasi bilan universitetga murojaat qiladi.",
     size=13, color=MUTED, spacing=1.25)

cards = [
    ("208", "OTM", "108 davlat · 70 nodavlat · 30 xorijiy filial", BLUE),
    ("47.7%", "Qamrov", "18–23 yosh qamrovi — 2014 yilda 6.8% edi", MINT),
    ("HEMIS", "Yagona tizim", "2021/22 dan barcha OTM uchun majburiy", AMBER),
]
for i, (num, lab, sub, col) in enumerate(cards):
    l = ML + i * (3.83 + 0.22)
    card(s, l, 3.80, 3.83, 1.85, CARD)
    dot(s, l + 0.45, 4.10, 0.16, col, MSO_SHAPE.RECTANGLE)
    text(s, l + 0.45, 4.38, 3.0, 0.60, num, size=34, bold=True, color=WHITE)
    text(s, l + 0.45, 4.92, 3.0, 0.28, lab, size=12, bold=True, color=col)
    text(s, l + 0.45, 5.22, 2.95, 0.55, sub, size=11, color=MUTED, spacing=1.2)

text(s, ML, 5.95, CW, 0.45,
     [[("Rasmiy keys: ", {"bold": True, "color": MUTED}),
       ("«…ma’lumotlarni izlash, uzun hujjatlarni o‘rganish, asosiy mazmunni ajratib olish… "
        "ko‘p vaqt talab qilishi mumkin.»", {"color": WHITE})]], size=13, spacing=1.2)
source(s, "Manba: Milliy statistika qo‘mitasi (stat.uz, 2025) · Oliy ta’lim, fan va innovatsiyalar vazirligi (2025) "
          "· NEXUS30 EdTech keysi, 3-bo‘lim")
footer(s, 1)

# =====================================================================
# 03 — problem
# =====================================================================
s = new_slide()
eyebrow(s, 2, "Problem", RED)

card(s, ML, 0.90, CW, 1.85, CARD2)
text(s, 1.20, 1.12, 3.3, 1.35, "4 manba", size=60, bold=True, color=RED)
text(s, 4.55, 1.28, 7.7, 0.55, "bitta oddiy savolga javob topish uchun", size=22, bold=True, color=WHITE)
text(s, 4.55, 1.82, 7.7, 0.85,
     "E’lonlar guruhi · dekanat koridori · qog‘oz nizom · kurator telefoni.\n"
     "Ma’lumot bor — lekin tarqoq va rolga bo‘linmagan.", size=14, color=MUTED, spacing=1.3)

text(s, ML, 3.30, 6.0, 0.30, "NIMA TO‘SQINLIK QILADI", size=12, bold=True, color=DIM, tracking=1.6)
probs = [
    ("1", "Qidiruv vaqti", "Javob hujjatda bor — lekin qaysi bandda, qaysi tilda, qaysi versiyada?", RED),
    ("2", "Takroriy savollar", "«Qancha qoldi», «arizam qayerda», «o‘qituvchi keldimi» — kuniga o‘nlab marta.", AMBER),
    ("3", "Rol chegarasi yo‘q", "Yo hammaga hamma narsa ochiq, yo hech kimga hech narsa.", BLUE),
]
for i, (n, title, sub, col) in enumerate(probs):
    l = ML + i * (3.83 + 0.22)
    card(s, l, 3.80, 3.83, 2.30, CARD)
    d = dot(s, l + 0.45, 4.25, 0.55, col)
    text(s, l + 0.45, 4.38, 0.55, 0.30, n, size=18, bold=True, color=BG, align=PP_ALIGN.CENTER)
    text(s, l + 0.45, 5.05, 3.0, 0.35, title, size=17, bold=True, color=WHITE)
    text(s, l + 0.45, 5.45, 2.95, 0.60, sub, size=11.5, color=MUTED, spacing=1.25)

source(s, "Muammo ta’rifi: NEXUS30 EdTech keysi (hamkor: Yandex), 3-bo‘lim «Muammo tavsifi»")
footer(s, 2)

# =====================================================================
# 04 — solution
# =====================================================================
s = new_slide()
eyebrow(s, 3, "Solution", MINT)
headline(s, "Bitta agent, besh rakurs — har javob manba bilan")

sols = [
    ("Rol-asosli RAG", "Savol → hujjat bo‘lagi → manba chipi. Agent javobni o‘ylab topmaydi.",
     "→ Qidiruv vaqti", RED),
    ("10 vosita (tool calling)", "Kontrakt qoldig‘i, turniket, ariza statusi — jonli bazadan, bitta savolda.",
     "→ Takroriy savollar", AMBER),
    ("Cheklov backendda", "Rol mos kelmasa vosita handleri umuman chaqirilmaydi — promptda emas, serverda.",
     "→ Rol chegarasi yo‘q", BLUE),
]
for i, (title, sub, answers, col) in enumerate(sols):
    l = ML + i * (3.83 + 0.22)
    card(s, l, 1.95, 3.83, 2.85, CARD)
    dot(s, l + 0.45, 2.30, 0.16, MINT, MSO_SHAPE.RECTANGLE)
    text(s, l + 0.45, 2.62, 3.0, 0.70, title, size=19, bold=True, color=WHITE, spacing=1.05)
    text(s, l + 0.45, 3.22, 2.95, 0.85, sub, size=12, color=MUTED, spacing=1.3)
    text(s, l + 0.45, 4.30, 2.95, 0.30, answers, size=11, bold=True, color=col)

card(s, ML, 5.10, CW, 1.10, CARD2)
text(s, 1.15, 5.32, 2.6, 0.70, "25 soniya", size=38, bold=True, color=MINT)
text(s, 4.10, 5.40, 8.1, 0.60,
     "butun demo ssenariysi — o‘lchangan sof tizim vaqti: 18 qadam, 5 rol, konsolda nol xato.",
     size=14, color=MUTED, spacing=1.25)
footer(s, 3)

# =====================================================================
# 05 — product
# =====================================================================
s = new_slide()
eyebrow(s, 4, "Product")
headline(s, "Savol → vosita → manba → javob", top=0.78, h=0.70)

steps = [
    ("1", "Savol", "Foydalanuvchi o‘z roli bilan kiradi"),
    ("2", "Vosita", "Registr ruxsatni serverda tekshiradi"),
    ("3", "Manba", "Hujjat + bo‘lim yoki turniket logi"),
    ("4", "Javob", "Faqat vosita natijasidan, disclaimer bilan"),
]
for i, (n, title, sub) in enumerate(steps):
    l = ML + i * (2.865 + 0.16)
    card(s, l, 1.70, 2.865, 1.00, CARD)
    text(s, l + 0.30, 1.86, 0.4, 0.30, n, size=13, bold=True, color=BLUE)
    text(s, l + 0.72, 1.84, 2.0, 0.30, title, size=14, bold=True, color=WHITE)
    text(s, l + 0.30, 2.20, 2.3, 0.40, sub, size=10.5, color=MUTED, spacing=1.15)

sw = 5.20
sl1, sl2 = 1.24, 1.24 + sw + 0.45
text(s, sl1, 2.92, sw, 0.28, "TALABA · manba chiplari va disclaimer", size=10, bold=True, color=DIM, tracking=1.0)
text(s, sl2, 2.92, sw, 0.28, "DEKANAT · «dars xavf ostida» — hisoblangan holat", size=10, bold=True, color=DIM,
     tracking=1.0)
shot(s, "01_chat_talaba.png", sl1, 3.22, sw)
shot(s, "02_dekanat_oqituvchilar.png", sl2, 3.22, sw)
footer(s, 4)

# =====================================================================
# 06 — why now
# =====================================================================
s = new_slide()
eyebrow(s, 5, "Why now", AMBER)
headline(s, "Uch shart bir vaqtda bajarildi")

whys = [
    ("280×", "LLM arzonlashdi", "Bir xil sifatdagi so‘rov narxi 2022-11 dan 2024-10 gacha shuncha tushdi.",
     "Stanford HAI, AI Index 2025", AMBER),
    ("5.8×", "Foydalanuvchi ko‘paydi", "Talaba soni 10 yilda: 264 ming → 1.5 mln. Xodim soni bunchalik o‘smadi.",
     "Milliy statistika qo‘mitasi, 2025", MINT),
    ("208/208", "Ma’lumot raqamli", "Barcha OTM HEMIS’da — integratsiya har universitet uchun emas, bir marta.",
     "VM 610-son (2020), 913-son (2024)", BLUE),
]
for i, (num, title, sub, src, col) in enumerate(whys):
    l = ML + i * (3.83 + 0.22)
    card(s, l, 1.95, 3.83, 3.10, CARD)
    text(s, l + 0.45, 2.25, 3.2, 0.90, num, size=44, bold=True, color=col)
    text(s, l + 0.45, 3.20, 3.0, 0.35, title, size=17, bold=True, color=WHITE)
    text(s, l + 0.45, 3.62, 2.95, 0.85, sub, size=11.5, color=MUTED, spacing=1.25)
    text(s, l + 0.45, 4.60, 2.95, 0.30, src, size=9, color=DIM)

card(s, ML, 5.35, CW, 0.85, CARD2)
text(s, 1.15, 5.58, 11.0, 0.45,
     "Model arzon, foydalanuvchi ko‘p, ma’lumot allaqachon raqamli — uchtasi ham ikki yil oldin yo‘q edi.",
     size=15, bold=True, color=WHITE)
footer(s, 5)

# =====================================================================
# 07 — use-case
# =====================================================================
s = new_slide()
eyebrow(s, 6, "Use-case")
headline(s, "Koridor o‘rniga bitta ekran", top=0.78, h=0.70)

col_w = 5.85
for i, (tag, q, before, after, shot_name, crop) in enumerate([
    ("TALABA", "«Arizam qayerda?»",
     "Bugun: dekanatga bor, navbatda tur, «ertaga kel».",
     "UniAgent: status va kim/qachon tarixi — bitta ekranda.",
     "04_talaba_arizalar.png", (0.0, 0.0, 0.0, 0.400)),
    ("TYUTOR / DEKANAT", "«Kim qarzdor, kim kelmadi?»",
     "Bugun: ro‘yxatni qo‘lda yig‘, muammoni kechqurun bil.",
     "UniAgent: svod tayyor, ogohlantirish dars davomida.",
     "03_tyutor_guruh.png", (0.0, 0.0, 0.26, 0.557)),
]):
    l = ML + i * (col_w + 0.23)
    card(s, l, 1.65, col_w, 4.35, CARD)
    text(s, l + 0.40, 1.88, 4.8, 0.28, tag, size=10, bold=True, color=BLUE, tracking=1.4)
    text(s, l + 0.40, 2.18, 5.1, 0.40, q, size=20, bold=True, color=WHITE)
    text(s, l + 0.40, 2.68, 5.1, 0.32, [[("Bugun: ", {"bold": True, "color": RED}), (before[7:], {"color": MUTED})]],
         size=12)
    text(s, l + 0.40, 3.02, 5.1, 0.32,
         [[("UniAgent: ", {"bold": True, "color": MINT}), (after[10:], {"color": MUTED})]], size=12)
    shot(s, shot_name, l + 0.40, 3.50, col_w - 0.80, crop)

card(s, ML, 6.20, CW, 0.62, CARD2)
text(s, 1.15, 6.38, 11.2, 0.32,
     [[("O‘lchangan tizim vaqti:  ", {"bold": True, "color": WHITE}),
       ("ariza yuborish 1.9 s · chek tasdiqlash 2.5 s · ariza tasdiqlash 3.6 s · agent javobi 0.5 s",
        {"color": MUTED})]], size=12)
footer(s, 6)

# =====================================================================
# 08 — business model
# =====================================================================
s = new_slide()
eyebrow(s, 7, "Business model", MINT)
headline(s, "SaaS litsenziya — to‘lovchi OTM, talaba emas")

for i, (num, title, sub, col) in enumerate([
    ("$9 000", "Yillik litsenziya", "Platforma, 5 rol, 10 vosita, RAG qidiruvi, yangilanish va qo‘llab-quvvatlash.",
     MINT),
    ("$6 000", "Bir martalik joriy etish", "HEMIS/SSO, Click yoki Payme, turniket va E-IMZO integratsiyasi.", BLUE),
]):
    l = ML + i * (5.85 + 0.23)
    card(s, l, 1.95, 5.85, 1.95, CARD)
    text(s, l + 0.45, 2.20, 3.0, 0.75, num, size=42, bold=True, color=col)
    text(s, l + 0.45, 3.00, 5.0, 0.32, title, size=16, bold=True, color=WHITE)
    text(s, l + 0.45, 3.38, 5.0, 0.40, sub, size=11.5, color=MUTED, spacing=1.2)

card(s, ML, 4.15, CW, 1.05, CARD2)
text(s, 1.15, 4.38, 3.4, 0.60, "$15 000", size=34, bold=True, color=WHITE)
text(s, 4.40, 4.48, 7.9, 0.45, "bitta OTM uchun birinchi yil qiymati — ikkinchi yildan $9 000 takrorlanuvchi.",
     size=14, color=MUTED)

card(s, ML, 5.40, CW, 0.90, CARD, AMBER)
text(s, 1.15, 5.60, 11.1, 0.55,
     [[("Narx — taxmin. ", {"bold": True, "color": AMBER}),
       ("O‘zbekistonda universitet SaaS litsenziyasi bo‘yicha ochiq benchmark topilmadi; "
        "yakuniy narx pilot natijasi bilan aniqlanadi.", {"color": MUTED})]], size=12, spacing=1.25)
footer(s, 7)

# =====================================================================
# 09 — market
# =====================================================================
s = new_slide()
eyebrow(s, 8, "Market size")
headline(s, "Katta global to‘lqin — O‘zbekistondan boshlanadi", top=0.78, h=0.75)

for i, (tag, val, label, detail, col) in enumerate([
    ("TAM", "$8.3 mlrd", "Jahon bozori",
     "«AI in education» (2025) → $57.2 mlrd (2033), CAGR 25.9%", BLUE),
    ("SAM", "≈$3.1 mln / yil", "O‘zbekiston OTMlari",
     "208 OTM × $15 000 birinchi yil qiymati", MINT),
    ("SOM", "≈$450 ming / yil", "3 yillik maqsad",
     "30 OTM (≈14%) — bitta pilot fakultetdan boshlab", AMBER),
]):
    t = 1.85 + i * 1.55
    card(s, ML, t, CW, 1.35, CARD)
    dot(s, ML + 0.40, t + 0.53, 0.28, col, MSO_SHAPE.RECTANGLE)
    text(s, ML + 0.90, t + 0.28, 4.4, 0.45,
         [[(f"{tag}  ·  ", {"color": col}), (val, {"color": WHITE})]], size=20, bold=True)
    text(s, ML + 0.90, t + 0.78, 4.4, 0.35, label, size=13, color=MUTED)
    text(s, 5.60, t + 0.52, 6.9, 0.55, detail, size=13, color=MUTED, spacing=1.2)

source(s, "Manba: Grand View Research, AI in Education Market (2026) · Oliy ta’lim vazirligi — 208 OTM (2025) "
          "· UniAgent narx taxmini")
footer(s, 8)

# =====================================================================
# 10 — competition
# =====================================================================
s = new_slide()
eyebrow(s, 9, "Competition", AMBER)
headline(s, "Hech bir muqobil beshta shartni birdan bajarmaydi", top=0.78, h=0.75)

cols = ["UniAgent", "HEMIS / LMS", "ChatGPT", "Telegram-bot"]
rows = [
    ("Rol-asosli xavfsizlik (serverda)", ["✓", "✓", "✕", "Qisman"]),
    ("O‘z hujjatidan javob + manba", ["✓", "✕", "Qisman", "✕"]),
    ("AI agent — savol-javob", ["✓", "✕", "✓", "Qisman"]),
    ("O‘zbek tili (ko‘p tilli qidiruv)", ["✓", "✓", "Qisman", "✓"]),
    ("Turniket / to‘lov integratsiyasi", ["✓", "Qisman", "✕", "✕"]),
    ("Javobgarlik qoidalari", ["✓", "—", "✕", "✕"]),
    ("Tizimda amal bajarish", ["✓", "✓", "✕", "Qisman"]),
]
col_x = [4.50, 6.53, 8.56, 10.59]
col_w2 = 2.03
row_h = 0.50
top0 = 1.87

card(s, col_x[0], top0, col_w2 + 0.09, row_h * (len(rows) + 1) + 0.16, CARD)
for j, name in enumerate(cols):
    text(s, col_x[j], top0 + 0.10, col_w2, 0.40, name, size=13, bold=True,
         color=WHITE if j == 0 else MUTED, align=PP_ALIGN.CENTER)
for i, (label, vals) in enumerate(rows):
    t = top0 + 0.58 + i * row_h
    sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(t - 0.05), Inches(CW), Emu(9525))
    _paint(sep, LINE)
    text(s, ML, t, 3.65, 0.45, label, size=12.5, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    for j, v in enumerate(vals):
        col = MINT if v == "✓" else (RED if v == "✕" else (AMBER if v == "Qisman" else DIM))
        text(s, col_x[j], t, col_w2, 0.45, v, size=14 if len(v) == 1 else 11.5, bold=(j == 0),
             color=col, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

source(s, "«Javobgarlik qoidalari» = majburiy disclaimer, manba ko‘rsatish va xulosani fakt sifatida bermaslik "
          "— keys 9-bo‘limi talabi.", top=6.35)
footer(s, 9)

# =====================================================================
# 11 — traction
# =====================================================================
s = new_slide()
eyebrow(s, 10, "Traction", MINT)
headline(s, "Prototip emas — ishlaydigan tizim", top=0.78, h=0.70)

for i, (num, label, col) in enumerate([
    ("281", "avtomatik test — hammasi o‘tadi", MINT),
    ("10", "AI vositasi · 5 rol · 9 modul", BLUE),
    ("25 s", "o‘lchangan demo — 6 daqiqa budjetida", AMBER),
]):
    l = ML + i * (3.83 + 0.22)
    card(s, l, 1.62, 3.83, 1.42, CARD)
    text(s, l + 0.40, 1.78, 3.2, 0.65, num, size=36, bold=True, color=col)
    text(s, l + 0.40, 2.48, 3.1, 0.40, label, size=11.5, color=MUTED, spacing=1.2)

facts = [
    "43 demo foydalanuvchi · 4 guruh · 2 fakultet",
    "10 hujjatli ko‘p tilli korpus (uz/ru/en) → 28 indekslangan bo‘lak",
    "RAG sifati: top-1 moslik 15/23 → 20/23",
    "100% sintetik ma’lumot — keysning qat’iy talabi",
    "S0–S14: 15 sessiya, har biri DoD tekshiruvi bilan yopilgan",
]
for i, f in enumerate(facts):
    t = 3.35 + i * 0.42
    dot(s, ML + 0.02, t + 0.09, 0.11, MINT, MSO_SHAPE.RECTANGLE)
    text(s, ML + 0.32, t, 5.6, 0.35, f, size=12, color=MUTED)

text(s, 6.90, 3.28, 5.7, 0.28, "ADMIN PANELI · jonli demo statistikasi", size=10, bold=True, color=DIM, tracking=1.0)
shot(s, "05_admin_statistika.png", 6.90, 3.58, 5.73, (0.0, 0.0, 0.0, 0.410))

card(s, ML, 6.05, CW, 0.62, CARD2)
text(s, 1.15, 6.23, 11.2, 0.32,
     [[("Keys talablari:  ", {"bold": True, "color": WHITE}),
       ("rolga mos agent ✓ · summarizatsiya ✓ · axborot qidirish ✓ · kontekstga mos javob ✓ · qulay interfeys ✓",
        {"color": MINT})]], size=12)
footer(s, 10)

# =====================================================================
# 12 — roadmap
# =====================================================================
s = new_slide()
eyebrow(s, 11, "Roadmap")
headline(s, "Pilotdan universitetgacha — uch bosqich")

for i, (when, title, sub, col) in enumerate([
    ("1-oy", "Pilot fakultet", "Real hujjat korpusi, HEMIS/SSO sinxronizatsiyasi, foydalanish statistikasi.", MINT),
    ("2-oy", "Integratsiya", "Click yoki Payme webhook, real turniket, PostgreSQL + boshqariladigan vektor baza.",
     BLUE),
    ("3–6 oy", "Butun universitet", "E-IMZO bilan rasmiy hujjat aylanmasi, push va Telegram bildirishnoma kanali.",
     AMBER),
]):
    l = ML + i * (3.83 + 0.22)
    card(s, l, 1.95, 3.83, 2.60, CARD)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(1.95), Inches(3.83), Inches(0.075))
    _paint(bar, col)
    text(s, l + 0.45, 2.30, 3.0, 0.35, when, size=13, bold=True, color=col, tracking=1.2)
    text(s, l + 0.45, 2.72, 3.1, 0.45, title, size=20, bold=True, color=WHITE)
    text(s, l + 0.45, 3.30, 2.95, 0.95, sub, size=12, color=MUTED, spacing=1.3)

card(s, ML, 4.85, CW, 0.95, CARD2)
text(s, 1.15, 5.10, 11.2, 0.45,
     [[("Keyin: ", {"bold": True, "color": WHITE}), ("2-yil — 10 OTM  ·  3-yil — 30 OTM", {"color": MINT})]],
     size=16, bold=True)

text(s, ML, 6.10, CW, 0.50,
     "Har bosqich mavjud kod ustiga quriladi: modellar, endpointlar va RBAC o‘zgarmaydi — "
     "seed generatori o‘rniga integratsiya adapteri qo‘yiladi.", size=12, color=DIM, spacing=1.25)
footer(s, 11)

# =====================================================================
# 13 — ask
# =====================================================================
s = new_slide()
eyebrow(s, 12, "The ask", RED)
headline(s, "Bitta pilot fakultet — bir oyda")

card(s, ML, 1.95, 7.35, 4.10, CARD)
asks = [
    ("G‘alaba — EdTech treki", "NEXUS30 finalida, hamkor Yandex bilan", MINT),
    ("Pilot universitet", "Bitta fakultet, 1 oy, real hujjat korpusi bilan", BLUE),
    ("Yandex model resursi", "llm/client.py da bitta klass — qolgan kod o‘zgarmaydi", AMBER),
    ("Mentorlik", "HEMIS integratsiyasi va ma’lumot maxfiyligi bo‘yicha", RED),
]
for i, (title, sub, col) in enumerate(asks):
    t = 2.28 + i * 0.92
    dot(s, ML + 0.42, t + 0.13, 0.20, col, MSO_SHAPE.RECTANGLE)
    text(s, ML + 0.90, t, 6.2, 0.32, title, size=16, bold=True, color=WHITE)
    text(s, ML + 0.90, t + 0.34, 6.2, 0.32, sub, size=12, color=MUTED)

card(s, 8.48, 1.95, 4.15, 4.10, CARD2)
text(s, 8.88, 2.70, 3.4, 1.15, "1 fakultet", size=34, bold=True, color=WHITE)
text(s, 8.88, 3.40, 3.4, 0.60, "1 oy", size=34, bold=True, color=MINT)
sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.88), Inches(4.25), Inches(3.35), Emu(9525))
_paint(sep, LINE)
text(s, 8.88, 4.55, 3.35, 1.60,
     "Pilotdan chiqadigan natija: real foydalanish statistikasi, aniq litsenziya narxi va birinchi referens OTM.",
     size=12.5, color=MUTED, spacing=1.3)

text(s, ML, 6.25, CW, 0.35,
     "Tizim bugun to‘liq ishlaydi — pilot uchun yozilishi kerak bo‘lgan yangi modul yo‘q, faqat integratsiya adapterlari.",
     size=12, color=DIM)
footer(s, 12)

# =====================================================================
# 14 — team
# =====================================================================
s = new_slide()
eyebrow(s, 13, "Jamoa", MINT)
headline(s, "Jamoa")

card(s, ML, 1.95, 7.35, 3.30, CARD)
av = dot(s, ML + 0.45, 2.35, 1.05, BLUE)
text(s, ML + 0.45, 2.68, 1.05, 0.45, "MS", size=24, bold=True, color=BG, align=PP_ALIGN.CENTER)
text(s, ML + 1.80, 2.42, 5.2, 0.45, "Mannonov Sarabek", size=24, bold=True, color=WHITE)
text(s, ML + 1.80, 2.92, 5.2, 0.32, "Full-stack + AI · loyihaning yagona muallifi", size=13, bold=True, color=BLUE)
text(s, ML + 0.45, 3.90, 6.5, 1.00,
     "Backend (FastAPI, RBAC, RAG quvuri, agent registri), frontend (Next.js 16), 281 test — "
     "15 sessiyada, har biri DoD tekshiruvi bilan yopilgan.", size=13, color=MUTED, spacing=1.3)

card(s, 8.48, 1.95, 4.15, 3.30, CARD2)
text(s, 8.88, 2.30, 3.4, 0.32, "NEGA BU MUHIM", size=10, bold=True, color=DIM, tracking=1.4)
text(s, 8.88, 2.75, 3.35, 2.20,
     "Jamoa bir kishilik — shuning uchun arxitektura qoidalari qat’iy: RBAC bitta joyda, "
     "LLM chaqiruvi bitta modul orqali, biznes logika services’da. Yangi odam kirsa ham qoida buzilmaydi.",
     size=12.5, color=MUTED, spacing=1.35)

card(s, ML, 5.55, CW, 0.85, CARD)
text(s, 1.15, 5.80, 11.2, 0.40, "UniAgent  ·  NEXUS30 EdTech  ·  Mannonov Sarabek  ·  Avgust 2026",
     size=14, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, 13)

prs.save(str(OUT))
print("saved:", OUT, f"{OUT.stat().st_size / 1024:.0f} KB", "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
